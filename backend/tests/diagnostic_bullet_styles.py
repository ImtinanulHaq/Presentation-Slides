#!/usr/bin/env python
"""
DIAGNOSTIC: Check bullet style in PPTX generation
"""
import os
import sys
import django

os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'presentation_project.settings')
django.setup()

from presentation_app.models import Presentation
from presentation_app.pptx_generator import generate_pptx
from pptx import Presentation as PptxPresentation

print("=" * 80)
print("BULLET STYLE DIAGNOSTIC TEST")
print("=" * 80)

# Get a presentation with different bullet styles
presentations = Presentation.objects.filter(
    bullet_style__in=['numbered', 'bullet_elegant', 'bullet_modern', 'bullet_professional']
).order_by('-id')[:4]

if not presentations.exists():
    print("\n✗ No presentations with different bullet styles found in database")
    print("Please create some presentations with different bullet styles first")
    sys.exit(1)

print(f"\nFound {len(presentations)} presentations to test\n")

for pres in presentations:
    style = pres.bullet_style
    print(f"\nTesting: {pres.title}")
    print(f"Bullet Style in DB: '{style}'")
    print("-" * 80)
    
    try:
        # Generate PPTX with the presentation's bullet style
        pptx_file = generate_pptx(
            pres,
            template_name=pres.template or 'rose_elegance',
            slide_ratio=pres.slide_ratio or '16:9',
            bullet_style=style  # Pass the style from database
        )
        
        # Read the generated PPTX to check bullets
        from io import BytesIO
        pptx_bytes = BytesIO(pptx_file.getvalue())
        pptx = PptxPresentation(pptx_bytes)
        
        # Check slide 2 (usually has bullets)
        if len(pptx.slides) >= 2:
            slide = pptx.slides[1]
            bullets_found = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs[:3]:  # First 3 paragraphs
                        if para.text.strip():
                            bullets_found.append(para.text[:70])
            
            if bullets_found:
                print(f"✓ Generated PPTX successfully")
                print(f"  Bullets on slide 2:")
                for bullet in bullets_found:
                    print(f"    • {bullet}")
                
                # Check if the style matches
                first_bullet = bullets_found[0] if bullets_found else ""
                if style == 'numbered' and first_bullet.startswith('1.'):
                    print(f"  ✓ Style applied correctly: NUMBERED")
                elif style == 'bullet_elegant' and first_bullet.startswith('●'):
                    print(f"  ✓ Style applied correctly: ELEGANT")
                elif style == 'bullet_modern' and first_bullet.startswith('▸'):
                    print(f"  ✓ Style applied correctly: MODERN")
                elif style == 'bullet_professional' and first_bullet.startswith('■'):
                    print(f"  ✓ Style applied correctly: PROFESSIONAL")
                else:
                    print(f"  ⚠ Style may not be applied correctly")
                    print(f"    Expected: '{style}'")
                    print(f"    First bullet: '{first_bullet}'")
            else:
                print(f"⚠ No bullets found on slide 2")
        else:
            print(f"⚠ Presentation has fewer than 2 slides")
            
    except Exception as e:
        print(f"✗ Error: {type(e).__name__}: {e}")

print("\n" + "=" * 80)
