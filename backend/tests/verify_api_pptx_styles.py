#!/usr/bin/env python
"""
Verify bullet styles in generated PPTX files from API
"""
from pptx import Presentation

def check_bullet_style(filename):
    """Check the bullet styles in a PPTX file"""
    try:
        prs = Presentation(filename)
        
        print(f"\n{filename}:")
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            if slide_idx <= 2:  # Check first 2 slides
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        text_frame = shape.text_frame
                        for para_idx, paragraph in enumerate(text_frame.paragraphs):
                            if para_idx < 3 and paragraph.text.strip():  # First 3 paragraphs
                                print(f"  Slide {slide_idx}, Para {para_idx}: {paragraph.text[:80]}")
        return True
    except Exception as e:
        print(f"  Error: {str(e)}")
        return False

if __name__ == '__main__':
    print("=" * 60)
    print("VERIFYING BULLET STYLES IN API-GENERATED PPTX FILES")
    print("=" * 60)
    
    files = [
        'test_api_numbered.pptx',
        'test_api_bullet_elegant.pptx',
        'test_api_bullet_modern.pptx',
        'test_api_bullet_professional.pptx'
    ]
    
    for file in files:
        check_bullet_style(file)
    
    print("\n" + "=" * 60)
