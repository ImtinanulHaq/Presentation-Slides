from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
from .templates.template_registry import TemplateRegistry


def calculate_optimal_font_size(text, max_width_inches, base_size, min_size=12):
    """
    Calculate optimal font size based on text length
    Prevents text from overflowing slide boundaries
    
    Args:
        text: Text to measure
        max_width_inches: Maximum width available
        base_size: Starting font size
        min_size: Minimum allowed font size
    
    Returns:
        int: Optimal font size in points
    """
    if not text:
        return base_size
    
    # Character estimation: approximately 6 characters per inch at 54pt
    # More conservative to prevent overflow
    chars_per_inch = 6
    available_chars = max_width_inches * chars_per_inch
    text_length = len(text)
    
    if text_length > available_chars:
        # Calculate reduction factor and new font size
        reduction_factor = available_chars / text_length
        new_size = int(base_size * reduction_factor)
        return max(new_size, min_size)
    
    return base_size


def generate_pptx(presentation_obj, template_name=None, slide_ratio='16:9', bullet_style='numbered'):
    """
    Generate a PPTX file from a presentation object with template styling
    Applies Rose Elegance template with professional colors, fonts, gradients, and layout
    
    Args:
        presentation_obj: Presentation model instance
        template_name: Template name (default: 'rose_elegance')
        slide_ratio: Slide aspect ratio ('16:9', '1:1', or '2:3')
        bullet_style: Bullet style ('numbered', 'bullet_elegant', 'bullet_modern', 'bullet_professional')
    
    Returns:
        BytesIO: PPTX file in memory
    """
    import logging
    logger = logging.getLogger(__name__)
    logger.info(f"🎯 generate_pptx called with bullet_style='{bullet_style}'")
    
    # Load template
    template = TemplateRegistry.get_template(template_name)
    
    # Get template colors
    colors = template.colors
    fonts = template.fonts
    text_styling = template.text_styling
    bullets = template.bullets
    
    # Parse colors from hex to RGB
    def hex_to_rgb(hex_color):
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    
    # Rose Elegance Colors - Professional palette
    PRIMARY_BG = hex_to_rgb(colors['primary_background'])      # #9d8189 - Deep mauve background
    SECONDARY_BG = hex_to_rgb(colors['secondary_background'])  # #f4acb7 - Medium rose background
    PRIMARY_TEXT = hex_to_rgb(colors['primary_text'])          # #5B3A49 - Dark mauve headings
    SECONDARY_TEXT = hex_to_rgb(colors['secondary_text'])      # #6E4B57 - Rich brownish-mauve body text
    ACCENT_BULLET = hex_to_rgb(colors['accent_color_dark'])    # #5B3A49 - Dark mauve bullets
    ACCENT_LIGHT = hex_to_rgb(colors['accent_color_light'])    # #ffcad4 - Soft rose
    WHITE = hex_to_rgb(colors['white'])                        # #ffffff - White
    
    # Create presentation with selected slide ratio
    prs = Presentation()
    
    # Set slide dimensions based on ratio
    if slide_ratio == '1:1':
        # Square format: 7.5" × 7.5"
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(7.5)
    elif slide_ratio == '4:3':
        # 4:3 Legacy format: 10" × 7.5" (older projectors)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    elif slide_ratio == '2:3':
        # Portrait format: 6.67" × 10"
        prs.slide_width = Inches(6.67)
        prs.slide_height = Inches(10)
    else:  # Default to 16:9
        # Widescreen format: 10" × 5.625" (proper 16:9 aspect ratio)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
    
    # Get slides from presentation
    slides = presentation_obj.slides.all().order_by('slide_number')
    
    for slide_idx, slide_obj in enumerate(slides):
        # Choose layout based on slide type
        if slide_obj.slide_type == 'title':
            # Title slide layout with professional gradient background
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add gradient background: Primary (#9d8189) to Accent Light (#ffcad4)
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_angle = 45.0  # Diagonal gradient for visual interest
            fill.gradient_stops[0].color.rgb = PRIMARY_BG  # Dark mauve (#9d8189)
            fill.gradient_stops[1].color.rgb = ACCENT_LIGHT  # Soft rose (#ffcad4)
            
            # Add decorative accent bar at bottom
            bar_width = prs.slide_width
            bar_height_inches = prs.slide_height.inches * 0.12  # 12% of slide height
            bar_top = prs.slide_height.inches * 0.85  # 85% from top
            
            accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(bar_top), bar_width, Inches(bar_height_inches))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = SECONDARY_BG  # Medium rose accent
            accent_bar.line.color.rgb = SECONDARY_BG
            accent_bar.line.width = Pt(0)
            
            # Add title with better sizing and wrapping
            title_text = slide_obj.title or presentation_obj.title
            # More aggressive reduction: 32-36pt for title slides (prevent overflow)
            max_title_width = prs.slide_width.inches * 0.9
            title_font_size = calculate_optimal_font_size(title_text, max_title_width, 36, 24)
            
            # Title box with proper margins - centered for maximum impact
            title_left = prs.slide_width.inches * 0.05
            title_width = prs.slide_width.inches * 0.9
            title_top = prs.slide_height.inches * 0.25
            title_height = prs.slide_height.inches * 0.45
            
            title_box = slide.shapes.add_textbox(Inches(title_left), Inches(title_top), Inches(title_width), Inches(title_height))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = title_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(title_font_size)
            p.font.bold = True
            p.font.name = 'Poppins'  # Use Poppins Bold for headings
            p.font.color.rgb = PRIMARY_TEXT  # Warm cream (#ffe5d9)
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(16)
            p.line_spacing = 1.2
            
            # Add subtitle with professional styling
            if slide_obj.subtitle:
                subtitle_text = slide_obj.subtitle
                subtitle_font_size = calculate_optimal_font_size(subtitle_text, max_title_width, 20, 14)
                
                subtitle_top = prs.slide_height.inches * 0.65
                subtitle_box = slide.shapes.add_textbox(Inches(title_left), Inches(subtitle_top), Inches(title_width), Inches(prs.slide_height.inches * 0.2))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.word_wrap = True
                subtitle_frame.vertical_anchor = MSO_ANCHOR.TOP
                
                p = subtitle_frame.paragraphs[0]
                p.text = subtitle_text
                p.font.size = Pt(subtitle_font_size)
                p.font.name = 'Inter'  # Use Inter for body text
                p.font.color.rgb = SECONDARY_TEXT  # Soft sage (#d8e2dc)
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(0)
                p.line_spacing = 1.3
            
        elif slide_obj.slide_type == 'section':
            # Section divider slide with professional gradient background
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add gradient background: Secondary (#f4acb7) to Primary (#9d8189)
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_angle = 135.0  # Different angle for visual variety
            fill.gradient_stops[0].color.rgb = SECONDARY_BG  # Medium rose (#f4acb7)
            fill.gradient_stops[1].color.rgb = PRIMARY_BG  # Dark mauve (#9d8189)
            
            # Dynamic dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add title with professional sizing - centered vertically and horizontally
            title_text = slide_obj.title
            max_width = slide_width.inches * 0.9
            title_font_size = calculate_optimal_font_size(title_text, max_width, 36, 26)
            
            title_left = slide_width.inches * 0.05
            title_width = slide_width.inches * 0.9
            title_top = slide_height.inches * 0.3
            title_height = slide_height.inches * 0.4
            
            title_box = slide.shapes.add_textbox(Inches(title_left), Inches(title_top), Inches(title_width), Inches(title_height))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = title_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(title_font_size)
            p.font.bold = True
            p.font.name = 'Poppins'  # Poppins Bold for section titles
            p.font.color.rgb = PRIMARY_TEXT  # Warm cream (#ffe5d9)
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.25
            
        else:
            # Content slide layout with professional gradient background
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add gradient background: Primary (#9d8189) to Accent Light (#ffcad4)
            background = slide.background
            fill = background.fill
            fill.gradient()
            fill.gradient_angle = 45.0
            fill.gradient_stops[0].color.rgb = PRIMARY_BG  # Dark mauve (#9d8189)
            fill.gradient_stops[1].color.rgb = ACCENT_LIGHT  # Soft rose (#ffcad4)
            
            # Dynamic dimensions based on slide width
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            title_bar_height = Inches(0.85)
            
            # Add title bar with secondary color
            title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), slide_width, title_bar_height)
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = SECONDARY_BG  # Medium rose (#f4acb7)
            title_bar.line.color.rgb = SECONDARY_BG
            title_bar.line.width = Pt(0)
            
            # Add title text with professional styling
            title_text = slide_obj.title
            title_font_size = calculate_optimal_font_size(title_text, slide_width.inches * 0.85, 26, 18)
            
            title_frame = title_bar.text_frame
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            title_frame.margin_left = Inches(0.5)
            title_frame.margin_right = Inches(0.5)
            title_frame.margin_top = Inches(0.15)
            title_frame.margin_bottom = Inches(0.15)
            title_frame.word_wrap = True
            
            p = title_frame.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(title_font_size)
            p.font.bold = True
            p.font.name = 'Poppins'  # Use Poppins Bold for headings
            p.font.color.rgb = PRIMARY_TEXT  # Warm cream (#ffe5d9)
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.15
            p.space_after = Pt(0)
            
            # Add content area with professional spacing
            # Content box positioned below title bar
            content_left = Inches(0.5)
            content_top = Inches(1.1)
            content_width = slide_width - Inches(1.0)
            content_height = slide_height - Inches(1.6)
            
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            # Professional bullet handling with style support
            if slide_obj.bullets:
                bullets_list = slide_obj.bullets
                
                # Ensure bullets is a list
                if isinstance(bullets_list, list):
                    bullets_list = [b for b in bullets_list if b and str(b).strip() and str(b).lower() not in ['none', 'null']]
                else:
                    bullets_list = []
                
                # Limit to 6 bullets for readability
                bullets_list = bullets_list[:6]
                
                # Determine bullet symbol based on bullet_style
                # The LLM should have already formatted bullets with the correct symbol
                # This is a fallback in case the LLM doesn't follow instructions
                bullet_symbol = None
                if bullet_style:
                    style = str(bullet_style).strip().lower()
                    if style == 'bullet_elegant':
                        bullet_symbol = '●'
                    elif style == 'bullet_modern':
                        bullet_symbol = '▸'
                    elif style == 'bullet_professional':
                        bullet_symbol = '■'
                
                # Calculate optimal font size for all bullets on this slide based on count and content length
                # This ensures all bullets on a slide have the SAME size and fit within bounds
                num_bullets = len(bullets_list)
                
                # Smart sizing: More bullets = smaller font, but stay readable
                if num_bullets == 6:
                    # Max bullets (6): use 15pt for balanced appearance
                    optimal_bullet_font_size = 15
                    space_after_bullets = Pt(5)  # Compact spacing
                elif num_bullets >= 5:
                    # Crowded: reduce to 16pt
                    optimal_bullet_font_size = 16
                    space_after_bullets = Pt(6)  # Reduced spacing
                elif avg_bullet_length > 80:
                    # Long text bullets even if few: use 16pt
                    optimal_bullet_font_size = 16
                    space_after_bullets = Pt(6)
                else:
                    # Normal case: use 18pt
                    optimal_bullet_font_size = 18
                    space_after_bullets = Pt(8)  # Standard spacing
                
                # Render each bullet with professional styling
                para_idx = 0
                for bullet in bullets_list:
                    if not bullet or not str(bullet).strip():
                        continue
                    
                    # Handle bullet text
                    if isinstance(bullet, dict):
                        bullet_text = str(bullet.get('text', ''))
                        level = int(bullet.get('level', 0))
                    else:
                        bullet_text = str(bullet).strip()
                        level = 0
                    
                    # Clean bullet text (remove any existing bullet symbols or numbers)
                    bullet_text = bullet_text.lstrip('•▸▪◆■★✓→ -0123456789.)')
                    bullet_text = bullet_text.strip()
                    
                    # Skip if empty after cleaning
                    if not bullet_text:
                        continue
                    
                    # Add paragraph
                    if para_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # Format bullet based on selected style
                    if bullet_symbol is not None:
                        # Use the selected bullet symbol (●, ▸, or ■)
                        formatted_bullet = f"{bullet_symbol} {bullet_text}"
                    else:
                        # Use numbered format (1., 2., 3., etc.)
                        formatted_bullet = f"{para_idx + 1}. {bullet_text}"
                    
                    # Set paragraph properties
                    p.text = formatted_bullet
                    p.level = min(level, 1)  # Max 2 levels
                    
                    para_idx += 1
                    # Apply optimized bullet size (same for all bullets on this slide)
                    bullet_font_size = optimal_bullet_font_size
                    
                    # Use template's bullet color (not hardcoded)
                    # Get bullet color from template bullets config
                    template_bullet_color_key = bullets.get('bullet_color', 'secondary_text')
                    template_bullet_color_hex = colors.get(template_bullet_color_key, colors.get('secondary_text', '#6E4B57'))
                    bullet_color = hex_to_rgb(template_bullet_color_hex)  # Use template's bullet color
                    
                    # Apply professional styling with optimized size and template color
                    p.font.size = Pt(bullet_font_size)
                    p.font.color.rgb = bullet_color  # Template's bullet color (dark blue/navy for Warm Spectrum)
                    p.font.name = 'Inter'  # Inter Regular
                    p.font.bold = False
                    
                    # Professional spacing for readability - adjusted based on bullet density
                    p.space_before = Pt(2)  # Minimal space before bullet
                    p.space_after = space_after_bullets  # Dynamic spacing based on bullet count
                    p.line_spacing = 1.4  # Professional line spacing
            
            # Add speaker notes if present
            if slide_obj.speaker_notes:
                try:
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = str(slide_obj.speaker_notes)[:500]  # Limit notes length
                except:
                    pass  # Ignore notes errors
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
